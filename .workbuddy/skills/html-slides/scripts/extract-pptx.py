"""
extract-pptx.py — Extract content from PowerPoint files for HTML Slides conversion.

Usage:
    python extract-pptx.py <input.pptx> <output_dir>

Output structure:
    <output_dir>/slides.json      — Slide content (titles, text, notes, layout hints)
    <output_dir>/assets/          — Extracted images

Requires: python-pptx (pip install python-pptx)
"""

import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def extract_images(slide, output_assets_dir):
    """Extract images from a slide and save to assets directory."""
    images = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture type
            image = shape.image
            image_bytes = image.blob
            ext = image.content_type.split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            filename = f"slide_{slide.slide_id}_{shape.shape_id}.{ext}"
            filepath = output_assets_dir / filename
            with open(filepath, 'wb') as f:
                f.write(image_bytes)
            images.append({
                'filename': filename,
                'path': str(filepath),
                'width': shape.width,
                'height': shape.height,
                'left': shape.left,
                'top': shape.top,
            })
    return images


def extract_text_content(shape):
    """Extract text content from a shape, preserving hierarchy."""
    if not shape.has_text_frame:
        return None

    paragraphs = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level = para.level
        paragraphs.append({
            'text': text,
            'level': level,
            'is_bullet': level > 0,
        })

    return paragraphs


def extract_slide_data(pptx_path, output_dir):
    """Extract all content from a PowerPoint file."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    assets_dir = output_dir / 'assets'
    assets_dir.mkdir(exist_ok=True)

    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_info = {
            'index': i,
            'slide_id': slide.slide_id,
            'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown',
            'title': '',
            'content': [],
            'notes': '',
            'images': extract_images(slide, assets_dir),
        }

        # Extract title
        if slide.shapes.title:
            slide_info['title'] = slide.shapes.title.text.strip()

        # Extract all shapes content
        for shape in slide.shapes:
            text_content = extract_text_content(shape)
            if text_content:
                shape_type = 'title' if shape == slide.shapes.title else 'body'
                slide_info['content'].append({
                    'shape_type': shape_type,
                    'shape_name': shape.name,
                    'paragraphs': text_content,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                })

        # Extract notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            slide_info['notes'] = notes_text

        slides_data.append(slide_info)

    # Write output
    output_file = output_dir / 'slides.json'
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump({
            'source': str(pptx_path),
            'total_slides': len(slides_data),
            'slide_width': str(prs.slide_width),
            'slide_height': str(prs.slide_height),
            'slides': slides_data,
        }, f, ensure_ascii=False, indent=2)

    print(f"Extracted {len(slides_data)} slides to {output_dir}")
    print(f"  Content: {output_file}")
    print(f"  Images:  {assets_dir}/ ({sum(len(s['images']) for s in slides_data)} images)")

    return slides_data


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python extract-pptx.py <input.pptx> <output_dir>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2]

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    extract_slide_data(pptx_path, output_dir)
